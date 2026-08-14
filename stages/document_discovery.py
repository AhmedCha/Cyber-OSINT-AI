#!/usr/bin/env python3
import argparse
import logging
import os
import re
import tempfile
import urllib.parse
import urllib.request
from pathlib import Path
from typing import Any, Dict, List, Set, Tuple

import openpyxl
import pypdf
import docx
import pptx

from lib.common import setup_logging, slugify_company
from lib.config import load_env_file
from lib.docker_runner import run_docker_tool
from lib.email_patterns import is_infrastructure_hostname
from lib.json_utils import load_json, save_json
from lib.search import generate_search_query_variants
from lib.apify_utils import run_apify_actor
from lib.db import get_db_connection, upsert_records

logger = logging.getLogger(__name__)

DEFAULT_FILE_TYPES = ["pdf", "doc", "docx", "xls", "xlsx", "ppt", "pptx"]
APIFY_GOOGLE_SEARCH_ACTOR_ID = "nFJndFXA5zjCTuudP"


def run_metagoofil_for_domain(domain: str, output_dir: Path) -> List[Dict[str, Any]]:
    """
    Runs Metagoofil via Docker for a given domain, downloads discovered files,
    and returns a list of metadata entries.
    """
    logger.info(f"[{domain}] Running Metagoofil document discovery...")
    results: List[Dict[str, Any]] = []

    domain_out_dir = output_dir / "metagoofil" / domain
    domain_out_dir.mkdir(parents=True, exist_ok=True)

    with tempfile.TemporaryDirectory() as tmp_dir:
        tmp_path = Path(tmp_dir)
        file_types_str = ",".join(DEFAULT_FILE_TYPES)

        cmd_args = [
            "-v",
            f"{tmp_path}:/data",
            "metagoofil",
            "-d",
            domain,
            "-t",
            file_types_str,
            "-w",
            "-o",
            "/data",
            "-n",
            "25",
            "-l",
            "50",
        ]

        try:
            run_docker_tool(tool_name="metagoofil", extra_args=cmd_args, timeout=600)
        except Exception as e:
            logger.error(f"[{domain}] Metagoofil run failed or timed out: {e}")

        for downloaded_file in tmp_path.glob("*"):
            if downloaded_file.is_file():
                dest_file = domain_out_dir / downloaded_file.name
                dest_file.write_bytes(downloaded_file.read_bytes())
                results.append(
                    {
                        "filename": downloaded_file.name,
                        "filepath": str(dest_file),
                        "domain": domain,
                        "source": "metagoofil",
                        "url": None,
                        "extracted_metadata": {},
                    }
                )

    return results


def run_apify_document_search(
    company_name: str,
    domains: List[str],
    existing_filenames: Set[str],
    output_dir: Path,
) -> List[Dict[str, Any]]:
    """
    Uses an Apify Google Search actor to find public document URLs across
    the web and for specific target domains.
    """
    logger.info(f"Running Apify web-wide document discovery for '{company_name}'...")
    discovered_files: List[Dict[str, Any]] = []

    query_variants = generate_search_query_variants(company_name)
    queries_list = [f'"{var}"' for var in query_variants]

    for domain in domains:
        if not is_infrastructure_hostname(domain):
            queries_list.append(f"site:{domain}")

    queries = "\n".join(queries_list)

    run_input = {
        "queries": queries,
        "maxPagesPerQuery": 3,
        "countryCode": "us",
        "searchLanguage": "en",
        "mobileResults": False,
        "fileTypes": DEFAULT_FILE_TYPES,
        "includeUnfilteredResults": True,
        "saveHtml": False,
    }

    items = run_apify_actor(APIFY_GOOGLE_SEARCH_ACTOR_ID, run_input)

    for item in items:
        organic_results = item.get("organicResults", [])

        for res in organic_results:
            url = res.get("url")
            if not url:
                continue

            parsed_url = urllib.parse.urlparse(url)
            doc_domain = parsed_url.netloc or "web"
            filename = urllib.parse.unquote(os.path.basename(parsed_url.path))

            unsafe_chars = ["/", "\\", "\x00", os.sep]
            if os.altsep:
                unsafe_chars.append(os.altsep)
            for char in unsafe_chars:
                filename = filename.replace(char, "_")

            ext = os.path.splitext(filename)[1].lstrip(".").lower()
            if not filename or ext not in DEFAULT_FILE_TYPES:
                continue

            if filename in existing_filenames:
                continue

            domain_out_dir = output_dir / "apify" / doc_domain
            domain_out_dir.mkdir(parents=True, exist_ok=True)
            dest_file = domain_out_dir / filename

            try:
                req = urllib.request.Request(
                    url,
                    headers={
                        "User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"
                    },
                )
                with urllib.request.urlopen(req, timeout=15) as response:
                    content = response.read()
                    content_type = response.headers.get("Content-Type", "").lower()

                    is_valid = False
                    if ext == "pdf" and content.startswith(b"%PDF"):
                        is_valid = True
                    elif ext in ["docx", "pptx", "xlsx"] and content.startswith(b"PK"):
                        is_valid = True
                    elif any(
                        valid_type in content_type
                        for valid_type in [
                            "pdf",
                            "document",
                            "msword",
                            "presentation",
                            "powerpoint",
                            "sheet",
                            "excel",
                        ]
                    ):
                        is_valid = True

                    if not is_valid:
                        logger.warning(
                            f"[{doc_domain}] Content verification failed for {url} (Ext: {ext}, Content-Type: {content_type}). Skipping save."
                        )
                        continue

                    dest_file.write_bytes(content)

                existing_filenames.add(filename)
                discovered_files.append(
                    {
                        "filename": filename,
                        "filepath": str(dest_file),
                        "domain": doc_domain,
                        "source": "apify-google",
                        "url": url,
                        "extracted_metadata": {
                            "title": res.get("title"),
                        },
                    }
                )
                logger.debug(
                    f"[{doc_domain}] Successfully downloaded document: {filename}"
                )
            except Exception as e:
                logger.warning(
                    f"[{doc_domain}] Failed to download Apify result {url}: {e}"
                )

    return discovered_files


def extract_pdf_content(file_path: Path) -> Tuple[str, Dict[str, Any]]:
    text = ""
    metadata = {}
    try:
        reader = pypdf.PdfReader(file_path)
        if reader.metadata:
            metadata = {
                "author": str(reader.metadata.author or ""),
                "title": str(reader.metadata.title or ""),
                "creator": str(reader.metadata.creator or ""),
                "producer": str(reader.metadata.producer or ""),
            }
        for page in reader.pages:
            page_text = page.extract_text()
            if page_text:
                text += page_text + " "
    except Exception as e:
        logger.warning(f"Failed to extract PDF content from {file_path.name}: {e}")
    return text, metadata


def extract_docx_content(file_path: Path) -> Tuple[str, Dict[str, Any]]:
    text = ""
    metadata = {}
    try:
        doc = docx.Document(str(file_path))
        cp = doc.core_properties
        metadata = {
            "author": cp.author or "",
            "title": cp.title or "",
            "subject": cp.subject or "",
            "created": str(cp.created) if cp.created else "",
        }
        text += " ".join([p.text for p in doc.paragraphs])
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += " " + cell.text
    except Exception as e:
        logger.warning(f"Failed to extract DOCX content from {file_path.name}: {e}")
    return text, metadata


def extract_pptx_content(file_path: Path) -> Tuple[str, Dict[str, Any]]:
    text = ""
    metadata = {}
    try:
        prs = pptx.Presentation(str(file_path))
        cp = prs.core_properties
        metadata = {
            "author": cp.author or "",
            "title": cp.title or "",
            "subject": cp.subject or "",
        }
        for slide in prs.slides:
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    text_frame = getattr(shape, "text_frame", None)
                    if text_frame:
                        text += " " + getattr(text_frame, "text", "")
    except Exception as e:
        logger.warning(f"Failed to extract PPTX content from {file_path.name}: {e}")
    return text, metadata


def extract_xlsx_content(file_path: Path) -> Tuple[str, Dict[str, Any]]:
    text = ""
    metadata = {}
    try:
        wb = openpyxl.load_workbook(file_path, data_only=True)
        if wb.properties:
            metadata = {
                "creator": wb.properties.creator or "",
                "title": wb.properties.title or "",
            }
        text_bits = []
        for sheet in wb.worksheets:
            for row in sheet.iter_rows(values_only=True):
                for cell in row:
                    if cell is not None:
                        text_bits.append(str(cell))
        text = " ".join(text_bits)
    except Exception as e:
        logger.warning(f"Failed to extract XLSX content from {file_path.name}: {e}")
    return text, metadata


def verify_file_content(
    file_path_str: str, company_name: str, domain: str
) -> Tuple[bool, Dict[str, Any]]:
    path = Path(file_path_str)
    ext = path.suffix.lower()

    text = ""
    metadata = {}

    if ext == ".pdf":
        text, metadata = extract_pdf_content(path)
    elif ext == ".docx":
        text, metadata = extract_docx_content(path)
    elif ext == ".pptx":
        text, metadata = extract_pptx_content(path)
    elif ext == ".xlsx":
        text, metadata = extract_xlsx_content(path)
    else:
        logger.info(
            f"Skipping content extraction for unsupported/legacy file format: {ext}"
        )
        return False, metadata

    text_lower = text.lower()
    search_terms = {company_name.lower(), domain.lower()}

    words = [w.lower() for w in re.split(r"\W+", company_name) if len(w) >= 3]
    search_terms.update(words)

    verified = any(term in text_lower for term in search_terms if term)
    return verified, metadata


def parse_arguments() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="OSINT Stage: Document Discovery & Verification"
    )
    parser.add_argument(
        "--company",
        required=True,
        help="Target company name",
    )
    return parser.parse_args()


def print_summary(results: List[Dict[str, Any]]) -> None:
    total = len(results)
    metagoofil_count = sum(1 for r in results if r["discovery_method"] == "metagoofil")
    apify_count = sum(1 for r in results if r["discovery_method"] == "apify-google")
    verified_count = sum(1 for r in results if r["content_verified"])

    print("\n" + "=" * 60)
    print("               DOCUMENT DISCOVERY SUMMARY")
    print("=" * 60)
    print(f"Total Files Found:             {total}")
    print(f"  - Metagoofil Discovery:      {metagoofil_count}")
    print(f"  - Apify Google Discovery:    {apify_count}")
    print(f"Content Verified Files:        {verified_count} / {total}")
    print("=" * 60 + "\n")


def main() -> None:
    setup_logging()
    logging.getLogger("pypdf").setLevel(logging.CRITICAL)
    load_env_file()

    args = parse_arguments()

    company_slug = slugify_company(args.company)
    company_dir = Path("output") / company_slug
    domains_file = company_dir / "domains.json"

    if not domains_file.exists():
        logger.error(
            f"Input file not found: {domains_file}. Please run domain discovery first."
        )
        return

    domains_data = load_json(domains_file)
    if not domains_data:
        logger.warning("No domains found in input file.")
        return

    domains = [item["domain"] for item in domains_data if "domain" in item]

    all_document_records: List[Dict[str, Any]] = []
    seen_filenames: Set[str] = set()

    for domain in domains:
        if is_infrastructure_hostname(domain):
            logger.info(
                f"Skipping infrastructure domain for Metagoofil discovery: {domain}"
            )
            continue

        mg_results = run_metagoofil_for_domain(domain, company_dir)
        for doc in mg_results:
            seen_filenames.add(doc["filename"])
            all_document_records.append(doc)

    ap_results = run_apify_document_search(
        args.company, domains, seen_filenames, company_dir
    )
    for doc in ap_results:
        all_document_records.append(doc)

    final_output: List[Dict[str, Any]] = []
    for doc in all_document_records:
        verified, parsed_meta = verify_file_content(
            doc["filepath"], args.company, doc["domain"]
        )

        merged_meta = doc.get("extracted_metadata", {})
        merged_meta.update(parsed_meta)

        final_output.append(
            {
                "filename": doc["filename"],
                "source_domain": doc["domain"],
                "discovery_method": doc["source"],
                "url": doc.get("url"),
                "extracted_metadata": merged_meta,
                "content_verified": verified,
            }
        )

    output_file = company_dir / "documents.json"
    save_json(output_file, final_output)

    # Save to db

    try:
        with get_db_connection() as conn:
            upsert_records(
                conn, "raw_documents", company_slug, final_output, "filename"
            )
    except Exception as e:
        logger.warning(f"Database sync failed for raw_documents: {e}")

    print_summary(final_output)


if __name__ == "__main__":
    main()
