"""
Documents category: text extraction from the formats document_discovery.py
downloads (pdf/docx/pptx/xlsx), and the one-LLM-call-per-document usability
summary pass. Unlike the other categories, this isn't a run_filter_pass
batch - each document gets its own call since the full extracted text (up
to --max-doc-chars) is the payload.
"""

from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

from lib.llm.client import ollama_chat_json
from lib.llm.config import BASE_SYSTEM_PROMPT, SUPPORTED_DOC_EXTENSIONS

USABILITY_LEVELS = {"high", "medium", "low", "unusable"}

DOC_SUMMARY_SCHEMA = {
    "type": "object",
    "properties": {
        "summary": {"type": "string"},
        "usability": {"type": "string", "enum": sorted(USABILITY_LEVELS)},
        "mentions_target_company": {"type": "boolean"},
    },
    "required": ["summary", "usability", "mentions_target_company"],
}


def extract_text(path: Path) -> Tuple[str, Optional[str]]:
    """Returns (text, error). Only handles the formats document_discovery.py
    downloads (pdf/docx/pptx/xlsx)."""
    ext = path.suffix.lower()
    try:
        if ext == ".pdf":
            import pypdf

            reader = pypdf.PdfReader(path)
            return " ".join(p.extract_text() or "" for p in reader.pages), None
        if ext == ".docx":
            import docx

            doc = docx.Document(str(path))
            parts = [p.text for p in doc.paragraphs]
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        parts.append(cell.text)
            return " ".join(parts), None
        if ext == ".pptx":
            import pptx

            prs = pptx.Presentation(str(path))
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False):
                        tf = getattr(shape, "text_frame", None)
                        if tf:
                            parts.append(getattr(tf, "text", ""))
            return " ".join(parts), None
        if ext == ".xlsx":
            import openpyxl

            wb = openpyxl.load_workbook(path, data_only=True)
            parts = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    for cell in row:
                        if cell is not None:
                            parts.append(str(cell))
            return " ".join(parts), None
        return "", f"Unsupported extension for text extraction: {ext}"
    except Exception as e:
        return "", f"Extraction failed: {e}"


def truncate_text(text: str, max_chars: int) -> Tuple[str, bool]:
    text = " ".join(text.split())  # normalize whitespace
    if len(text) <= max_chars:
        return text, False
    return text[:max_chars], True


def summarize_document(
    doc: Dict[str, Any],
    company_name: str,
    host: str,
    model: str,
    max_chars: int,
    temperature: float,
    timeout: int,
    warnings: List[str],
    debug: bool = False,
    num_gpu: Optional[int] = None,
    num_ctx: Optional[int] = None,
    num_predict: Optional[int] = None,
) -> Dict[str, Any]:
    base = {
        "filename": doc["filename"],
        "source_domain": doc["source_domain"],
        "filepath": doc["filepath"],
        "url": doc.get("url"),
        "file_exists": doc["file_exists"],
        "content_verified": doc.get("content_verified", False),
        "summary": None,
        "usability": "unusable",
        "mentions_target_company": None,
        "truncated": False,
        "error": None,
    }

    if not doc.get("file_exists"):
        base["error"] = "File not found on disk - skipped LLM summarization."
        return base

    path = Path(doc["filepath"])
    if path.suffix.lower() not in SUPPORTED_DOC_EXTENSIONS:
        base["error"] = f"Unsupported file type for extraction: {path.suffix}"
        return base

    text, extract_err = extract_text(path)
    if extract_err:
        base["error"] = extract_err
        return base
    if not text.strip():
        base["error"] = "No extractable text (empty or scanned/image-only document)."
        return base

    truncated_text, was_truncated = truncate_text(text, max_chars)
    base["truncated"] = was_truncated

    user_prompt = f"""Target company: {company_name}
Document filename: {doc["filename"]}
Document text{" (TRUNCATED - do not assume anything about content beyond what follows)" if was_truncated else ""}:
---
{truncated_text}
---

Based ONLY on the text above, respond with JSON:
{{
  "summary": "<one sentence, plain language, describing what this document actually is and whether it looks useful for an OSINT report on the target company>",
  "usability": "<one of: high, medium, low, unusable>",
  "mentions_target_company": <true/false - does the text explicitly reference the target company or its name/domain?>
}}"""

    result, error = ollama_chat_json(
        host,
        model,
        BASE_SYSTEM_PROMPT,
        user_prompt,
        DOC_SUMMARY_SCHEMA,
        temperature,
        timeout=timeout,
        debug=debug,
        debug_label=f"document:{doc['filename']}",
        num_gpu=num_gpu,
        num_ctx=num_ctx,
        num_predict=num_predict,
    )

    if result is None:
        base["error"] = f"LLM summarization failed: {error}"
        warnings.append(f"[documents] '{doc['filename']}': {base['error']}")
        return base

    summary = result.get("summary")
    usability = result.get("usability")
    mentions = result.get("mentions_target_company")

    base["summary"] = summary if isinstance(summary, str) else None
    base["usability"] = usability if usability in USABILITY_LEVELS else "unusable"
    base["mentions_target_company"] = mentions if isinstance(mentions, bool) else None

    if base["summary"] is None:
        warnings.append(
            f"[documents] '{doc['filename']}': model returned no usable summary field."
        )

    return base
