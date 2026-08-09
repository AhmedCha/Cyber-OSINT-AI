#!/usr/bin/env python3
"""
OSINT Stage: LLM Filter

Takes the company's aggregate_results.json (produced by aggregate_results.py)
and runs it through a local LLM (Ollama, llama3.1:8b by default - see
install_tools.sh) to:

  1. Read each discovered document's extracted text and produce a
     one-sentence usability summary (for later report inclusion decisions).
  2. Strip noise from the structured OSINT data (duplicate/irrelevant
     SpiderFoot findings, generic ISP/infra metadata, low-value pattern-
     generated emails, etc.) while keeping everything transparent - nothing
     is silently deleted, it is labeled "excluded" with a reason.

DESIGN PRINCIPLE - the LLM never owns the schema or the facts:
  - The LLM is only ever asked to return a verdict (keep/exclude + a short
    note) about a record that ALREADY exists in the aggregate data, using
    that record's own identifier (email / domain / employee id).
  - Every verdict is grounded against the original input set after parsing.
    Any identifier the model invents that doesn't match an input record is
    dropped and logged as a hallucination - it can never appear in the
    output.
  - If the model omits a record entirely, or fails after retries, the
    pipeline "fails open": the original record is kept, unmodified, with a
    warning attached - we never let an LLM hiccup silently delete real
    OSINT data.
  - The Python code always emits the same fixed top-level structure
    (company / domains / emails / employees / breaches / documents /
    dns_infra / warnings / stats / model), regardless of which model
    produced the content. Swapping llama3.1:8b for another model changes
    the *quality* of summaries/verdicts, never the *shape* of the output.

PERFORMANCE / LOW-RESOURCE HARDWARE NOTES:
  - List-based passes (domains/emails/employees/breaches) are sent in small
    batches (--batch-size, default 8) rather than one giant call, since a
    single call with dozens of records can exceed the model's context
    window (Ollama's default num_ctx is commonly 4096) and take a very long
    time to generate on CPU/iGPU-only hardware, risking timeouts.
  - --debug streams tokens live to the terminal as the model generates them
    (Ollama's stream=true API), so you can see it's actually working instead
    of staring at a blank terminal until a timeout.
  - --num-gpu lets you force layers onto CPU (0 = fully CPU) per request,
    useful for testing whether GPU offload is the source of instability
    (e.g. display blackouts) without touching your global Ollama service
    config.

Usage:
    python -m stages.llm_filter --company "Resys Consultants"
    python -m stages.llm_filter --company "Resys Consultants" --debug --batch-size 5 --timeout 900
    python -m stages.llm_filter --company "Resys Consultants" --num-gpu 0   # force CPU-only, for stability testing
"""

import argparse
import json
import logging
import os
import sys
import time
import urllib.error
import urllib.request
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Callable, Dict, List, Optional, Tuple

from lib.common import setup_logging, slugify_company
from lib.json_utils import load_json, save_json

try:
    from lib.config import load_env_file
except ImportError:  # pragma: no cover - keep this stage runnable standalone

    def load_env_file() -> None:
        return None


logger = logging.getLogger(__name__)

# =====================================================================
# CONFIG
# =====================================================================

DEFAULT_OLLAMA_HOST = os.environ.get("OLLAMA_HOST", "http://localhost:11434")
DEFAULT_MODEL = os.environ.get("OLLAMA_MODEL", "llama3.1:8b")
DEFAULT_TIMEOUT = 600  # generous default for CPU/iGPU-only hardware
DEFAULT_MAX_RETRIES = 3
DEFAULT_MAX_DOC_CHARS = (
    4000  # kept conservative relative to a likely 4096-token context window
)
DEFAULT_BATCH_SIZE = 8  # records per LLM call for list-based filter passes

SUPPORTED_DOC_EXTENSIONS = {".pdf", ".docx", ".pptx", ".xlsx"}

USABILITY_LEVELS = {"high", "medium", "low", "unusable"}
EMAIL_TIERS = {"confirmed", "likely", "speculative", "noise"}


# =====================================================================
# SYSTEM PROMPT (shared anti-hallucination contract)
# =====================================================================

BASE_SYSTEM_PROMPT = """You are a precise OSINT data-triage assistant. You are given data that has \
ALREADY been collected by automated reconnaissance tools (domain enumeration, \
email pattern generation, LinkedIn scraping, breach lookups, document \
discovery). Your only job is to judge and summarize what you are given.

STRICT RULES - violating any of these makes your output useless and unsafe:
1. NEVER invent, guess, assume, or infer any fact that is not explicitly \
present in the data given to you in this message. If something is unknown, \
say so or leave it out - do not fill gaps with plausible-sounding guesses.
2. NEVER fabricate new identifiers (emails, domains, names, filenames, URLs). \
You may only refer to identifiers that appear verbatim in the input below.
3. NEVER change, "correct", normalize, or reformat identifiers (emails, \
domains, filenames) - copy them exactly as given.
4. Output ONLY valid JSON matching the requested schema. No markdown code \
fences, no prose before or after the JSON, no explanations outside the JSON \
fields themselves.
5. If you are not confident about a judgment, prefer the more conservative \
label (e.g. "low" usability, "speculative" tier, keep=true) rather than \
discarding data - a human will review your output before anything is \
published in a report.
6. You have no internet access and cannot verify anything beyond the text \
given to you in this message.
"""


# =====================================================================
# OLLAMA CLIENT
# =====================================================================


def check_ollama_available(host: str, timeout: int = 10) -> bool:
    try:
        req = urllib.request.Request(f"{host.rstrip('/')}/api/tags")
        with urllib.request.urlopen(req, timeout=timeout) as resp:
            return resp.status == 200
    except Exception as e:
        logger.error(f"Ollama is not reachable at {host}: {e}")
        return False


def _build_options(
    temperature: float,
    num_gpu: Optional[int],
    num_ctx: Optional[int],
    num_predict: Optional[int],
) -> Dict[str, Any]:
    options: Dict[str, Any] = {"temperature": temperature, "seed": 42}
    if num_gpu is not None:
        options["num_gpu"] = num_gpu
    if num_ctx is not None:
        options["num_ctx"] = num_ctx
    if num_predict is not None:
        options["num_predict"] = num_predict
    return options


def _post_ollama_chat(
    host: str,
    model: str,
    system_prompt: str,
    user_prompt: str,
    schema: Optional[Dict[str, Any]],
    options: Dict[str, Any],
    timeout: int,
    use_schema_format: bool,
    debug: bool = False,
    debug_label: str = "",
) -> str:
    """Low-level call to Ollama's /api/chat. Returns the full message content
    string. When debug=True, streams the response and prints tokens live to
    stdout as they arrive (Ollama's stream=true NDJSON API)."""
    body: Dict[str, Any] = {
        "model": model,
        "messages": [
            {"role": "system", "content": system_prompt},
            {"role": "user", "content": user_prompt},
        ],
        "stream": debug,
        "options": options,
    }

    if schema is not None and use_schema_format:
        body["format"] = schema
    else:
        body["format"] = "json"

    data = json.dumps(body).encode("utf-8")
    req = urllib.request.Request(
        f"{host.rstrip('/')}/api/chat",
        data=data,
        headers={"Content-Type": "application/json"},
        method="POST",
    )

    if not debug:
        with urllib.request.urlopen(req, timeout=timeout) as resp:
            raw = json.loads(resp.read().decode("utf-8"))
        return raw.get("message", {}).get("content", "")

    # --- streaming / debug path -----------------------------------
    print(
        f"\n--- [{debug_label}] live model output "
        + "-" * max(0, 40 - len(debug_label))
    )
    full_content = []
    start = time.time()
    with urllib.request.urlopen(req, timeout=timeout) as resp:
        for line in resp:
            line = line.strip()
            if not line:
                continue
            try:
                chunk = json.loads(line.decode("utf-8"))
            except json.JSONDecodeError:
                continue
            piece = chunk.get("message", {}).get("content", "")
            if piece:
                print(piece, end="", flush=True)
                full_content.append(piece)
            if chunk.get("done"):
                elapsed = time.time() - start
                eval_count = chunk.get("eval_count")
                eval_duration_ns = chunk.get("eval_duration")
                rate = ""
                if eval_count and eval_duration_ns:
                    tok_per_sec = eval_count / (eval_duration_ns / 1e9)
                    rate = f" | {eval_count} tokens @ {tok_per_sec:.1f} tok/s"
                print(f"\n--- [{debug_label}] done in {elapsed:.1f}s{rate} ---\n")
    return "".join(full_content)


def ollama_chat_json(
    host: str,
    model: str,
    system_prompt: str,
    user_prompt: str,
    schema: Optional[Dict[str, Any]] = None,
    temperature: float = 0.0,
    timeout: int = DEFAULT_TIMEOUT,
    max_retries: int = DEFAULT_MAX_RETRIES,
    debug: bool = False,
    debug_label: str = "",
    num_gpu: Optional[int] = None,
    num_ctx: Optional[int] = None,
    num_predict: Optional[int] = None,
) -> Tuple[Optional[Dict[str, Any]], Optional[str]]:
    """Calls the local LLM and returns (parsed_json, error). Retries on
    malformed JSON by re-prompting; falls back to generic json mode if the
    schema-constrained call itself fails (older Ollama versions)."""
    use_schema_format = schema is not None
    last_error = None
    current_user_prompt = user_prompt
    options = _build_options(temperature, num_gpu, num_ctx, num_predict)

    for attempt in range(1, max_retries + 1):
        attempt_start = time.time()
        try:
            content = _post_ollama_chat(
                host,
                model,
                system_prompt,
                current_user_prompt,
                schema,
                options,
                timeout,
                use_schema_format,
                debug=debug,
                debug_label=f"{debug_label} attempt {attempt}",
            )
        except urllib.error.HTTPError as e:
            body = ""
            try:
                body = e.read().decode("utf-8", errors="replace")[:500]
            except Exception:
                pass
            if use_schema_format and e.code == 400:
                logger.warning(
                    "Ollama rejected structured `format` schema (HTTP 400). "
                    "Falling back to generic JSON mode."
                )
                use_schema_format = False
                continue
            last_error = f"HTTP {e.code}: {e.reason} | body: {body}"
            logger.warning(
                f"[{debug_label}] Ollama call failed (attempt {attempt}/{max_retries}, "
                f"{time.time() - attempt_start:.1f}s): {last_error}"
            )
        except urllib.error.URLError as e:
            last_error = f"Connection error: {e.reason}"
            logger.warning(
                f"[{debug_label}] Ollama call failed (attempt {attempt}/{max_retries}, "
                f"{time.time() - attempt_start:.1f}s): {last_error}"
            )
        except TimeoutError as e:
            last_error = f"Timed out after {timeout}s"
            logger.warning(
                f"[{debug_label}] Ollama call timed out (attempt {attempt}/{max_retries}). "
                f"Consider raising --timeout, lowering --batch-size, or trying --num-gpu 0 "
                f"if GPU offload seems unstable on this hardware."
            )
        except Exception as e:
            last_error = str(e)
            logger.warning(
                f"[{debug_label}] Ollama call failed (attempt {attempt}/{max_retries}, "
                f"{time.time() - attempt_start:.1f}s): {last_error}"
            )
        else:
            elapsed = time.time() - attempt_start
            try:
                parsed = json.loads(content)
                logger.info(f"[{debug_label}] completed in {elapsed:.1f}s")
                return parsed, None
            except json.JSONDecodeError as e:
                last_error = f"Invalid JSON from model: {e}"
                logger.warning(
                    f"[{debug_label}] Model returned invalid JSON after {elapsed:.1f}s "
                    f"(attempt {attempt}/{max_retries}). Retrying with correction."
                )
                current_user_prompt = (
                    user_prompt
                    + "\n\nYour previous response was not valid JSON. "
                    + "Respond with ONLY valid JSON matching the schema - no other text."
                )

        time.sleep(min(2 * attempt, 6))

    return None, last_error or "Unknown error"


# =====================================================================
# TEXT EXTRACTION (documents pass)
# =====================================================================


def extract_text(path: Path) -> Tuple[str, Optional[str]]:
    """Returns (text, error). Only handles the formats document_discovery.py
    downloads (pdf/docx/pptx/xlsx)."""
    ext = path.suffix.lower()
    try:
        if ext == ".pdf":
            import pypdf

            reader = pypdf.PdfReader(path)
            return " ".join(p.extract_text() or "" for p in reader.pages), None
        if ext == ".docx":
            import docx

            doc = docx.Document(str(path))
            parts = [p.text for p in doc.paragraphs]
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        parts.append(cell.text)
            return " ".join(parts), None
        if ext == ".pptx":
            import pptx

            prs = pptx.Presentation(str(path))
            parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if getattr(shape, "has_text_frame", False):
                        tf = getattr(shape, "text_frame", None)
                        if tf:
                            parts.append(getattr(tf, "text", ""))
            return " ".join(parts), None
        if ext == ".xlsx":
            import openpyxl

            wb = openpyxl.load_workbook(path, data_only=True)
            parts = []
            for sheet in wb.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    for cell in row:
                        if cell is not None:
                            parts.append(str(cell))
            return " ".join(parts), None
        return "", f"Unsupported extension for text extraction: {ext}"
    except Exception as e:
        return "", f"Extraction failed: {e}"


def truncate_text(text: str, max_chars: int) -> Tuple[str, bool]:
    text = " ".join(text.split())  # normalize whitespace
    if len(text) <= max_chars:
        return text, False
    return text[:max_chars], True


DOC_SUMMARY_SCHEMA = {
    "type": "object",
    "properties": {
        "summary": {"type": "string"},
        "usability": {"type": "string", "enum": sorted(USABILITY_LEVELS)},
        "mentions_target_company": {"type": "boolean"},
    },
    "required": ["summary", "usability", "mentions_target_company"],
}


def summarize_document(
    doc: Dict[str, Any],
    company_name: str,
    host: str,
    model: str,
    max_chars: int,
    temperature: float,
    timeout: int,
    warnings: List[str],
    debug: bool = False,
    num_gpu: Optional[int] = None,
    num_ctx: Optional[int] = None,
    num_predict: Optional[int] = None,
) -> Dict[str, Any]:
    base = {
        "filename": doc["filename"],
        "source_domain": doc["source_domain"],
        "filepath": doc["filepath"],
        "file_exists": doc["file_exists"],
        "content_verified": doc.get("content_verified", False),
        "summary": None,
        "usability": "unusable",
        "mentions_target_company": None,
        "truncated": False,
        "error": None,
    }

    if not doc.get("file_exists"):
        base["error"] = "File not found on disk - skipped LLM summarization."
        return base

    path = Path(doc["filepath"])
    if path.suffix.lower() not in SUPPORTED_DOC_EXTENSIONS:
        base["error"] = f"Unsupported file type for extraction: {path.suffix}"
        return base

    text, extract_err = extract_text(path)
    if extract_err:
        base["error"] = extract_err
        return base
    if not text.strip():
        base["error"] = "No extractable text (empty or scanned/image-only document)."
        return base

    truncated_text, was_truncated = truncate_text(text, max_chars)
    base["truncated"] = was_truncated

    user_prompt = f"""Target company: {company_name}
Document filename: {doc["filename"]}
Document text{" (TRUNCATED - do not assume anything about content beyond what follows)" if was_truncated else ""}:
---
{truncated_text}
---

Based ONLY on the text above, respond with JSON:
{{
  "summary": "<one sentence, plain language, describing what this document actually is and whether it looks useful for an OSINT report on the target company>",
  "usability": "<one of: high, medium, low, unusable>",
  "mentions_target_company": <true/false - does the text explicitly reference the target company or its name/domain?>
}}"""

    result, error = ollama_chat_json(
        host,
        model,
        BASE_SYSTEM_PROMPT,
        user_prompt,
        DOC_SUMMARY_SCHEMA,
        temperature,
        timeout=timeout,
        debug=debug,
        debug_label=f"document:{doc['filename']}",
        num_gpu=num_gpu,
        num_ctx=num_ctx,
        num_predict=num_predict,
    )

    if result is None:
        base["error"] = f"LLM summarization failed: {error}"
        warnings.append(f"[documents] '{doc['filename']}': {base['error']}")
        return base

    summary = result.get("summary")
    usability = result.get("usability")
    mentions = result.get("mentions_target_company")

    base["summary"] = summary if isinstance(summary, str) else None
    base["usability"] = usability if usability in USABILITY_LEVELS else "unusable"
    base["mentions_target_company"] = mentions if isinstance(mentions, bool) else None

    if base["summary"] is None:
        warnings.append(
            f"[documents] '{doc['filename']}': model returned no usable summary field."
        )

    return base


# =====================================================================
# GENERIC "FILTER PASS" (domains / emails / employees / breaches)
# =====================================================================


def _chunked(items: List[Any], size: int) -> List[List[Any]]:
    if size <= 0:
        return [items]
    return [items[i : i + size] for i in range(0, len(items), size)]


def run_filter_pass(
    category: str,
    records: List[Dict[str, Any]],
    identifier_field: str,
    compact_fn: Callable[[Dict[str, Any]], Dict[str, Any]],
    verdict_item_schema: Dict[str, Any],
    system_prompt: str,
    instructions: str,
    host: str,
    model: str,
    temperature: float,
    timeout: int,
    warnings: List[str],
    batch_size: int = DEFAULT_BATCH_SIZE,
    debug: bool = False,
    num_gpu: Optional[int] = None,
    num_ctx: Optional[int] = None,
    num_predict: Optional[int] = None,
) -> Tuple[List[Dict[str, Any]], List[Dict[str, Any]]]:
    """Sends compact versions of `records` to the LLM in batches of
    `batch_size`, gets back a keep/exclude verdict per identifier, grounds
    every verdict against the real input set, and merges the verdict onto
    the ORIGINAL (untouched) record. Fails open: any record the model
    doesn't mention is kept with a note.

    Returns (kept, excluded) - both lists contain full original records plus
    an added `_llm_verdict` block.
    """
    if not records:
        return [], []

    valid_ids = {str(r.get(identifier_field, "")).strip().lower() for r in records}

    schema = {
        "type": "object",
        "properties": {"verdicts": {"type": "array", "items": verdict_item_schema}},
        "required": ["verdicts"],
    }

    verdict_map: Dict[str, Dict[str, Any]] = {}
    batches = _chunked(records, batch_size)

    for batch_num, batch in enumerate(batches, 1):
        compact_records = [compact_fn(r) for r in batch]
        label = f"{category} batch {batch_num}/{len(batches)}"

        user_prompt = f"""{instructions}

Input records (JSON array, {len(compact_records)} items):
{json.dumps(compact_records, ensure_ascii=False, indent=2)}

Respond with JSON: {{"verdicts": [ ... one object per input record, using the exact identifier field value from the input above ... ]}}
You must return exactly one verdict object per input record. Do not add records that are not in the input above."""

        result, error = ollama_chat_json(
            host,
            model,
            system_prompt,
            user_prompt,
            schema,
            temperature,
            timeout=timeout,
            debug=debug,
            debug_label=label,
            num_gpu=num_gpu,
            num_ctx=num_ctx,
            num_predict=num_predict,
        )

        if result is None:
            warnings.append(
                f"[{category}] Batch {batch_num}/{len(batches)} LLM filter pass failed ({error}). "
                f"Keeping those {len(batch)} record(s) unfiltered."
            )
            continue

        raw_verdicts = result.get("verdicts", [])
        if not isinstance(raw_verdicts, list):
            warnings.append(
                f"[{category}] Batch {batch_num}/{len(batches)} response had no valid 'verdicts' array. "
                f"Keeping those {len(batch)} record(s) unfiltered."
            )
            continue

        for v in raw_verdicts:
            if not isinstance(v, dict):
                continue
            vid = str(v.get(identifier_field, "")).strip().lower()
            if not vid:
                # Either the model omitted the identifier field, or
                # identifier_field doesn't match the key name actually used
                # in the schema/compact_fn for this category - both are bugs
                # worth surfacing rather than silently dropping the verdict.
                warnings.append(
                    f"[{category}] Batch {batch_num}/{len(batches)}: verdict item missing "
                    f"'{identifier_field}' field (got keys: {sorted(v.keys())}). Dropped."
                )
                continue
            if vid not in valid_ids:
                warnings.append(
                    f"[{category}] Dropped hallucinated identifier not present in input: {v.get(identifier_field)!r}"
                )
                continue
            verdict_map[vid] = v

    kept: List[Dict[str, Any]] = []
    excluded: List[Dict[str, Any]] = []

    for record in records:
        rid = str(record.get(identifier_field, "")).strip().lower()
        verdict = verdict_map.get(rid)

        if verdict is None:
            annotated = dict(record)
            annotated["_llm_verdict"] = {
                "keep": True,
                "note": "not_evaluated_by_model_fail_open",
            }
            kept.append(annotated)
            continue

        keep = verdict.get("keep", True)
        if not isinstance(keep, bool):
            keep = True

        annotated = dict(record)
        annotated["_llm_verdict"] = {
            k: v for k, v in verdict.items() if k != identifier_field
        }
        (kept if keep else excluded).append(annotated)

    return kept, excluded


# =====================================================================
# CATEGORY-SPECIFIC COMPACTORS, PROMPTS & SCHEMAS
# =====================================================================


def compact_domain(d: Dict[str, Any]) -> Dict[str, Any]:
    return {
        "domain": d.get("domain"),
        "sources": d.get("sources", []),
        "dns_validated": d.get("dns_validated"),
    }


DOMAIN_VERDICT_ITEM_SCHEMA = {
    "type": "object",
    "properties": {
        "domain": {"type": "string"},
        "keep": {"type": "boolean"},
        "note": {"type": "string"},
    },
    "required": ["domain", "keep"],
}

DOMAIN_INSTRUCTIONS = (
    "You are reviewing candidate domains discovered by automated domain-enumeration "
    "tools (theHarvester, SpiderFoot, certificate transparency logs) for a target company. "
    "Mark keep=true for domains that plausibly belong to or are operated by the target "
    "organization. Mark keep=false ONLY for domains that are clearly unrelated false "
    "positives (e.g. an unrelated third-party service, a coincidental substring match, "
    "a parked/placeholder domain) - do not exclude a domain just because you are unsure."
)


def compact_email(e: Dict[str, Any]) -> Dict[str, Any]:
    return {
        "email": e.get("email"),
        "employee": e.get("employee"),
        "sources": e.get("sources", []),
        "validation_status": e.get("validation_status"),
        "confidence": e.get("confidence"),
        "is_catch_all": e.get("is_catch_all"),
    }


EMAIL_VERDICT_ITEM_SCHEMA = {
    "type": "object",
    "properties": {
        "email": {"type": "string"},
        "keep": {"type": "boolean"},
        "tier": {"type": "string", "enum": sorted(EMAIL_TIERS)},
        "note": {"type": "string"},
    },
    "required": ["email", "keep", "tier"],
}

EMAIL_INSTRUCTIONS = (
    "You are reviewing candidate email addresses for a target company. Assign a tier "
    "based PRIMARILY on validation_status and confidence, not on how plausible the "
    "naming pattern looks:\n"
    "- 'confirmed': validation_status='deliverable' (SMTP validation actually succeeded).\n"
    "- 'likely': validation_status='unknown' but confidence is GREATER THAN 0 (some "
    "partial validation signal exists, e.g. a catch-all domain).\n"
    "- 'speculative': confidence is 0 (SMTP validation was never run, was inconclusive, "
    "or failed) - use this even for a common, plausible-looking naming pattern like "
    "first.last@ or f.last@. A confidence of 0 means the address is UNVALIDATED, not "
    "that the pattern is unlikely, but you must not call it 'likely' - that would "
    "overstate certainty that validation never actually confirmed.\n"
    "- 'noise': generic role addresses with no named employee, or a duplicate pattern "
    "already covered by a better candidate for the same person.\n"
    "NEVER assign 'confirmed' or 'likely' when confidence is 0, regardless of how "
    "standard the naming pattern looks. Set keep=false ONLY for tier='noise' entries "
    "that add no value - when in doubt, keep=true."
)


def enforce_email_tier_rules(
    records: List[Dict[str, Any]], warnings: List[str]
) -> None:
    """Code-level backstop for the confidence-vs-tier rule, applied AFTER the
    model's verdicts are merged in. Prompt instructions alone aren't
    reliable enough to guarantee this - an 8B model (or a future swapped-in
    model) can still call an unvalidated, confidence=0 address 'likely'
    because the naming pattern looks plausible. This corrects that
    deterministically, in-place, on both kept and excluded records, and
    logs every correction for transparency."""
    for r in records:
        verdict = r.get("_llm_verdict")
        if not isinstance(verdict, dict):
            continue
        tier = verdict.get("tier")
        if tier not in ("confirmed", "likely"):
            continue

        confidence = r.get("confidence")
        validation_status = r.get("validation_status")

        if tier == "confirmed" and validation_status != "deliverable":
            verdict["tier"] = "speculative" if (confidence or 0) == 0 else "likely"
            warnings.append(
                f"[emails] Downgraded '{r.get('email')}' from tier=confirmed to "
                f"tier={verdict['tier']}: validation_status is '{validation_status}', not 'deliverable'."
            )
        elif tier == "likely" and (confidence or 0) == 0:
            verdict["tier"] = "speculative"
            warnings.append(
                f"[emails] Downgraded '{r.get('email')}' from tier=likely to tier=speculative: "
                f"confidence is 0 (unvalidated / SMTP inconclusive), regardless of naming pattern."
            )


def compact_employee(e: Dict[str, Any]) -> Dict[str, Any]:
    identifier = e.get("public_identifier") or e.get("name")
    current = e.get("current_position") or {}
    if not isinstance(current, dict):
        current = {"value": current}
    positions = _as_list(e.get("all_positions"))
    trimmed_positions = [
        {
            "title": p.get("title"),
            "company_name": p.get("company_name"),
            "start_date": p.get("start_date"),
            "end_date": p.get("end_date"),
        }
        for p in positions[:6]
        if isinstance(p, dict)
    ]
    about = str(e.get("about") or "")[:400]
    return {
        "identifier": identifier,
        "name": e.get("name"),
        "job_title": e.get("job_title"),
        "matched_domain": e.get("matched_domain"),
        "current_position": current,
        "recent_positions": trimmed_positions,
        "about_excerpt": about,
        "services_offered": _as_list(e.get("services_offered"))[:6],
    }


EMPLOYEE_VERDICT_ITEM_SCHEMA = {
    "type": "object",
    "properties": {
        "identifier": {"type": "string"},
        "keep": {"type": "boolean"},
        "connection_to_target": {"type": "string"},
        "key_facts": {"type": "array", "items": {"type": "string"}},
        "note": {"type": "string"},
    },
    "required": ["identifier", "keep"],
}

EMPLOYEE_INSTRUCTIONS = (
    "You are reviewing people discovered via LinkedIn search who are linked to the "
    "target company (current or former employees, interns, consultants). For each "
    "person, in 'connection_to_target' state in one short phrase how they relate to "
    "the target company based ONLY on their position history given (e.g. 'current "
    "employee', 'former intern 2019-2021', 'consultant/trainer'). In 'key_facts' list "
    "up to 4 short, factual bullet points about them drawn ONLY from the data given "
    "(role, notable skills/certifications, tenure) - no speculation about personality, "
    "seniority you cannot verify, or anything not in the input. Set keep=false ONLY if "
    "the person has no discernible connection to the target company in the data given."
)


def _as_list(x: Any) -> List[Any]:
    """Defensively normalize a value that's supposed to be a list but, given
    breach_lookup.py isn't written yet and different lookup services return
    different shapes, might arrive as a dict, a scalar, or missing entirely."""
    if isinstance(x, list):
        return x
    if x is None:
        return []
    return [x]


def _is_structured_breach_report(data: Any) -> bool:
    """Detects the Apify/XposedOrNot-style breach report shape (has
    breachNames/status/riskLabel etc.) vs a SpiderFoot raw-event dump or
    some other unknown shape."""
    return isinstance(data, dict) and (
        "breachNames" in data or "status" in data or "riskScore" in data
    )


def _extract_structured_breach_summary(
    data: Dict[str, Any], source: Optional[str]
) -> Dict[str, Any]:
    """Pulls the genuinely high-value fields out of a large Apify/XposedOrNot
    breach report, dropping the verbose nested 'analytics' tree (industry
    stat breakdowns, yearwise histograms, treemap data, etc.) that adds
    nothing for an OSINT report but would otherwise dominate the LLM's
    context budget on slow hardware."""
    breach_details = _as_list(data.get("breachDetails"))
    trimmed_details = []
    for bd in breach_details[:10]:
        if not isinstance(bd, dict):
            continue
        trimmed_details.append(
            {
                "breach": bd.get("breach"),
                "industry": bd.get("industry"),
                "xposed_date": bd.get("xposed_date"),
                "xposed_records": bd.get("xposed_records"),
                "xposed_data": bd.get("xposed_data"),
                "password_risk": bd.get("password_risk"),
                "verified": bd.get("verified"),
            }
        )
    return {
        "source": source,
        "kind": "structured_breach_report",
        "status": data.get("status"),
        "breach_count": data.get("breachCount"),
        "breach_names": data.get("breachNames"),
        "risk_label": data.get("riskLabel"),
        "risk_score": data.get("riskScore"),
        "paste_count": data.get("pasteCount"),
        "breach_details": trimmed_details,
    }


def compact_breach(b: Dict[str, Any]) -> Dict[str, Any]:
    entries = []
    for group in _as_list(b.get("breaches"))[:5]:
        if not isinstance(group, dict):
            # Unexpected shape (e.g. a bare string/number breach entry) -
            # surface it as-is rather than crashing or silently dropping it.
            entries.append({"source": None, "kind": "unknown", "value": group})
            continue

        source = group.get("source") or group.get("type")
        data = group.get("data")

        if isinstance(data, dict) and _is_structured_breach_report(data):
            # Real breach report (e.g. Apify/XposedOrNot) - extract the
            # high-value summary fields, drop the noisy analytics tree.
            entries.append(_extract_structured_breach_summary(data, source))
            continue

        data_items = _as_list(data)[:5]
        if not data_items:
            entries.append(
                {"source": source, "kind": "empty", "raw_type": group.get("type")}
            )
            continue

        for item in data_items:
            if isinstance(item, dict):
                # SpiderFoot-style {type, module, data} raw event, or some
                # other dict shape we don't specifically recognize.
                entries.append(
                    {
                        "source": source,
                        "kind": "raw_event",
                        "type": item.get("type", group.get("type")),
                        "module": item.get("module"),
                        "data": item["data"] if "data" in item else item,
                    }
                )
            else:
                entries.append(
                    {
                        "source": source,
                        "kind": "raw_event",
                        "type": group.get("type"),
                        "module": None,
                        "data": item,
                    }
                )

    return {
        "email": b.get("email"),
        "services": b.get("services", []),
        "findings": entries,
    }


BREACH_VERDICT_ITEM_SCHEMA = {
    "type": "object",
    "properties": {
        "email": {"type": "string"},
        "keep": {"type": "boolean"},
        "exposure_summary": {"type": "string"},
        "note": {"type": "string"},
    },
    "required": ["email", "keep"],
}

BREACH_INSTRUCTIONS = (
    "You are reviewing breach/exposure lookup results for target-company emails. Each "
    "email's 'findings' list can contain two very different kinds of entries:\n"
    "1. kind='structured_breach_report' - a real breach-database lookup result (e.g. "
    "from XposedOrNot/Apify) with fields like status, breach_count, breach_names, "
    "risk_label, risk_score, and breach_details (each with breach name, xposed_date, "
    "xposed_records, xposed_data, password_risk). If status is 'breached' or "
    "breach_names is non-empty, this IS a genuine exposure - keep=true, and write "
    "'exposure_summary' as a concise sentence naming the breaches, how many records/what "
    "kind of data was exposed, and the risk_label - using ONLY the fields given.\n"
    "2. kind='raw_event' with source='spiderfoot' (or similar) - usually just SpiderFoot "
    "confirming an email address exists (type='Email Address', module='SpiderFoot UI', "
    "data equal to the email itself), with no actual breach name/date attached. This is "
    "NOT a real exposure - keep=false with a note explaining it's just an existence check, "
    "UNLESS its 'data' field clearly names a real breach/leak source, in which case treat "
    "it like a genuine finding.\n"
    "If an email has both a genuine structured_breach_report AND spiderfoot noise, keep=true "
    "and summarize only the genuine findings - do not mention the noise in exposure_summary."
)


def compact_darkweb(d: Dict[str, Any]) -> Dict[str, Any]:
    mentions = _as_list(d.get("mentions"))
    return {
        "target_key": f"{d.get('target')}||{d.get('target_type')}",
        "target": d.get("target"),
        "target_type": d.get("target_type"),
        "modules_checked": d.get("modules_checked", []),
        "mentions": mentions[:10],
    }


DARKWEB_VERDICT_ITEM_SCHEMA = {
    "type": "object",
    "properties": {
        "target_key": {"type": "string"},
        "keep": {"type": "boolean"},
        "exposure_summary": {"type": "string"},
        "note": {"type": "string"},
    },
    "required": ["target_key", "keep"],
}

DARKWEB_INSTRUCTIONS = (
    "You are reviewing dark web scan results for the target company (its domain, its "
    "name, and named individuals linked to it). Each record has a 'mentions' list of "
    "actual hits found on onion search engines / dark web indexes. Some mentions may be "
    "irrelevant false positives (e.g. an unrelated page that happens to contain a common "
    "word from the target's name). Set keep=true if the mentions genuinely appear to "
    "reference the target company or person, and write 'exposure_summary' describing "
    "what was found using ONLY the data given. Set keep=false ONLY if the mentions are "
    "clearly unrelated false positives - when unsure, keep=true."
)


def compact_infrastructure(i: Dict[str, Any]) -> Dict[str, Any]:
    return {
        "target_domain": i.get("target_domain"),
        "theHarvester": i.get("theHarvester", {}),
        "SpiderFoot": i.get("SpiderFoot", []),
    }


INFRASTRUCTURE_VERDICT_ITEM_SCHEMA = {
    "type": "object",
    "properties": {
        "target_domain": {"type": "string"},
        "keep": {"type": "boolean"},
        "ips": {"type": "array", "items": {"type": "string"}},
        "asns": {"type": "array", "items": {"type": "string"}},
        "whois_data": {"type": "array", "items": {"type": "string"}},
        "banners_and_tech": {"type": "array", "items": {"type": "string"}},
        "note": {"type": "string"},
    },
    "required": ["target_domain", "keep"],
}

INFRASTRUCTURE_INSTRUCTIONS = (
    "You are reviewing raw infrastructure discovery data (from theHarvester and SpiderFoot) "
    "for a target domain. Extract genuinely useful structured findings like real IP addresses, "
    "ASN/network ownership info, notable WHOIS fields, and interesting webserver or tech-stack "
    "banners. Explicitly ignore purely internal bookkeeping events, duplicate/redundant entries, "
    "and empty fields. Set keep=true if you found actionable infrastructure data. Set keep=false "
    "ONLY if the data is entirely noise, empty, or uninformative."
)

# =====================================================================
# ASSEMBLY & OUTPUT
# =====================================================================


def build_output_template(
    company: str, company_slug: str, model_cfg: Dict[str, Any]
) -> Dict[str, Any]:
    """The one place that defines the full, fixed shape of llm_filtered.json.
    Every key here is always present in the final output, regardless of
    which model produced the content."""
    return {
        "company": company,
        "company_slug": company_slug,
        "generated_at": datetime.now(timezone.utc).isoformat(),
        "model": model_cfg,
        "domains": {"kept": [], "excluded": []},
        "emails": {"kept": [], "excluded": []},
        "employees": {"kept": [], "excluded": []},
        "breaches": {"kept": [], "excluded": []},
        "darkweb": {"kept": [], "excluded": []},
        "infrastructure_insights": {"kept": [], "excluded": []},
        "documents": [],
        "dns_infra": {},
        "warnings": [],
        "stats": {},
    }


def main() -> None:
    setup_logging()
    load_env_file()
    args = parse_arguments()

    if args.debug:
        logging.getLogger().setLevel(logging.DEBUG)

    company_slug = slugify_company(args.company)
    company_dir = Path("output") / company_slug
    input_file = company_dir / "aggregate_results.json"

    if not input_file.exists():
        logger.error(
            f"Input file not found: {input_file}. Run aggregate_results.py first."
        )
        raise SystemExit(1)

    aggregate = load_json(input_file)
    if not aggregate:
        logger.error(f"{input_file} is empty or invalid.")
        raise SystemExit(1)

    if not check_ollama_available(args.ollama_host):
        logger.error(
            f"Ollama is not reachable at {args.ollama_host}. "
            "Start it (see install_tools.sh) before running this stage."
        )
        raise SystemExit(1)

    model_cfg = {"provider": "ollama", "name": args.model, "host": args.ollama_host}
    output = build_output_template(args.company, company_slug, model_cfg)
    warnings: List[str] = []

    common_kwargs = dict(
        host=args.ollama_host,
        model=args.model,
        temperature=args.temperature,
        timeout=args.timeout,
        warnings=warnings,
        batch_size=args.batch_size,
        debug=args.debug,
        num_gpu=args.num_gpu,
        num_ctx=args.num_ctx,
        num_predict=args.num_predict,
    )

    # --- Domains -----------------------------------------------------
    logger.info("Filtering domains...")
    kept, excluded = run_filter_pass(
        "domains",
        aggregate.get("domains", []),
        "domain",
        compact_domain,
        DOMAIN_VERDICT_ITEM_SCHEMA,
        BASE_SYSTEM_PROMPT,
        DOMAIN_INSTRUCTIONS,
        **common_kwargs,
    )
    output["domains"] = {"kept": kept, "excluded": excluded}

    # --- Infrastructure Insights ---------------------------------------
    logger.info("Filtering infrastructure raw data...")
    infra_raw_dict = aggregate.get("infrastructure_raw", {})
    infra_records = []

    for domain, data in infra_raw_dict.items():
        rec = dict(data)
        rec["target_domain"] = domain
        infra_records.append(rec)

    infra_kept, infra_excluded = run_filter_pass(
        "infrastructure",
        infra_records,
        "target_domain",
        compact_infrastructure,
        INFRASTRUCTURE_VERDICT_ITEM_SCHEMA,
        BASE_SYSTEM_PROMPT,
        INFRASTRUCTURE_INSTRUCTIONS,
        **common_kwargs,
    )
    output["infrastructure_insights"] = {"kept": infra_kept, "excluded": infra_excluded}

    # --- Emails --------------------------------------------------------
    logger.info("Filtering emails...")
    kept, excluded = run_filter_pass(
        "emails",
        aggregate.get("emails", []),
        "email",
        compact_email,
        EMAIL_VERDICT_ITEM_SCHEMA,
        BASE_SYSTEM_PROMPT,
        EMAIL_INSTRUCTIONS,
        **common_kwargs,
    )
    output["emails"] = {"kept": kept, "excluded": excluded}
    enforce_email_tier_rules(kept, warnings)
    enforce_email_tier_rules(excluded, warnings)

    # --- Employees -------------------------------------------------
    logger.info("Filtering employees...")

    def employee_identifier(e: Dict[str, Any]) -> str:
        return e.get("public_identifier") or e.get("name") or ""

    employees_with_id = []
    for e in aggregate.get("employees", []):
        e2 = dict(e)
        e2["public_identifier"] = employee_identifier(e)
        # "identifier" (not "public_identifier") is the key name used by both
        # compact_employee() and EMPLOYEE_VERDICT_ITEM_SCHEMA - it must match
        # exactly what's passed as identifier_field below, or every verdict
        # silently fails to ground and the whole pass fails open.
        e2["identifier"] = e2["public_identifier"]
        employees_with_id.append(e2)

    kept, excluded = run_filter_pass(
        "employees",
        employees_with_id,
        "identifier",
        compact_employee,
        EMPLOYEE_VERDICT_ITEM_SCHEMA,
        BASE_SYSTEM_PROMPT,
        EMPLOYEE_INSTRUCTIONS,
        **common_kwargs,
    )
    output["employees"] = {"kept": kept, "excluded": excluded}

    # --- Breaches --------------------------------------------------
    logger.info("Filtering breaches...")
    kept, excluded = run_filter_pass(
        "breaches",
        aggregate.get("breaches", []),
        "email",
        compact_breach,
        BREACH_VERDICT_ITEM_SCHEMA,
        BASE_SYSTEM_PROMPT,
        BREACH_INSTRUCTIONS,
        **common_kwargs,
    )
    output["breaches"] = {"kept": kept, "excluded": excluded}

    # --- Dark web ----------------------------------------------------
    logger.info("Filtering dark web scan results...")
    darkweb_records = []
    for d in aggregate.get("darkweb", []):
        d2 = dict(d)
        d2["target_key"] = f"{d.get('target')}||{d.get('target_type')}"
        darkweb_records.append(d2)

    with_mentions = [d for d in darkweb_records if d.get("mentions")]
    without_mentions = [d for d in darkweb_records if not d.get("mentions")]

    dw_kept, dw_excluded = run_filter_pass(
        "darkweb",
        with_mentions,
        "target_key",
        compact_darkweb,
        DARKWEB_VERDICT_ITEM_SCHEMA,
        BASE_SYSTEM_PROMPT,
        DARKWEB_INSTRUCTIONS,
        **common_kwargs,
    )
    # Targets with zero mentions need no LLM judgment - there's nothing to
    # filter, so they're kept automatically without spending a call on them.
    for d in without_mentions:
        annotated = dict(d)
        annotated["_llm_verdict"] = {"keep": True, "note": "no_dark_web_mentions_found"}
        dw_kept.append(annotated)
    output["darkweb"] = {"kept": dw_kept, "excluded": dw_excluded}

    # --- Documents (one LLM call per document) ----------------------
    documents = aggregate.get("documents", [])
    logger.info(f"Summarizing {len(documents)} document(s)...")
    doc_summaries = []
    for i, doc in enumerate(documents, 1):
        logger.info(f"  [{i}/{len(documents)}] {doc.get('filename')}")
        doc_summaries.append(
            summarize_document(
                doc,
                args.company,
                args.ollama_host,
                args.model,
                args.max_doc_chars,
                args.temperature,
                args.timeout,
                warnings,
                debug=args.debug,
                num_gpu=args.num_gpu,
                num_ctx=args.num_ctx,
                num_predict=args.num_predict,
            )
        )
    output["documents"] = doc_summaries

    # --- DNS infra: passthrough (no LLM value yet, kept for future noise-stripping) ---
    output["dns_infra"] = aggregate.get("dns_infra", {})

    output["warnings"] = warnings
    output["stats"] = {
        "domains_kept": len(output["domains"]["kept"]),
        "domains_excluded": len(output["domains"]["excluded"]),
        "emails_kept": len(output["emails"]["kept"]),
        "emails_excluded": len(output["emails"]["excluded"]),
        "employees_kept": len(output["employees"]["kept"]),
        "employees_excluded": len(output["employees"]["excluded"]),
        "breaches_kept": len(output["breaches"]["kept"]),
        "breaches_excluded": len(output["breaches"]["excluded"]),
        "darkweb_kept": len(output["darkweb"]["kept"]),
        "darkweb_excluded": len(output["darkweb"]["excluded"]),
        "darkweb_with_mentions": len(with_mentions),
        "infrastructure_kept": len(output["infrastructure_insights"]["kept"]),
        "infrastructure_excluded": len(output["infrastructure_insights"]["excluded"]),
        "documents_summarized": sum(1 for d in doc_summaries if d["summary"]),
        "documents_errored": sum(1 for d in doc_summaries if d["error"]),
        "warning_count": len(warnings),
    }

    output_file = company_dir / "llm_filtered.json"
    save_json(output_file, output, indent=2)

    print_summary(output["stats"], warnings, output_file)


def parse_arguments() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="OSINT Stage: LLM-based noise filtering & document summarization"
    )
    parser.add_argument("--company", required=True, help="Target company name")
    parser.add_argument("--model", default=DEFAULT_MODEL, help="Ollama model tag")
    parser.add_argument(
        "--ollama-host", default=DEFAULT_OLLAMA_HOST, help="Ollama API base URL"
    )
    parser.add_argument(
        "--temperature",
        type=float,
        default=0.0,
        help="Sampling temperature (0 = deterministic)",
    )
    parser.add_argument(
        "--max-doc-chars",
        type=int,
        default=DEFAULT_MAX_DOC_CHARS,
        help="Max characters of extracted document text sent to the model per document",
    )
    parser.add_argument(
        "--timeout",
        type=int,
        default=DEFAULT_TIMEOUT,
        help="Per-request timeout in seconds (raise this on slow/CPU-only hardware)",
    )
    parser.add_argument(
        "--batch-size",
        type=int,
        default=DEFAULT_BATCH_SIZE,
        help="Records per LLM call for domains/emails/employees/breaches passes "
        "(lower this if calls are timing out)",
    )
    parser.add_argument(
        "--debug",
        action="store_true",
        help="Stream the model's raw output live to the terminal as it generates, "
        "with per-call timing/throughput stats",
    )
    parser.add_argument(
        "--num-gpu",
        type=int,
        default=None,
        help="Ollama num_gpu option: number of layers to offload to GPU for this request "
        "(0 = force fully CPU-only, useful to test if GPU offload is unstable on this "
        "machine). Omit to use the Ollama service's default.",
    )
    parser.add_argument(
        "--num-ctx",
        type=int,
        default=None,
        help="Ollama num_ctx option: context window size for this request. Omit to use "
        "the model's default (commonly 4096).",
    )
    parser.add_argument(
        "--num-predict",
        type=int,
        default=None,
        help="Ollama num_predict option: caps max tokens generated per response, guards "
        "against runaway generation on slow hardware.",
    )
    return parser.parse_args()


def print_summary(
    stats: Dict[str, int], warnings: List[str], output_file: Path
) -> None:
    print("\n" + "=" * 60)
    print("               LLM FILTER SUMMARY")
    print("=" * 60)
    for label, value in stats.items():
        print(f"{label:<28}: {value}")
    if warnings:
        print("-" * 60)
        print(f"Warnings ({len(warnings)}):")
        for w in warnings[:10]:
            print(f"  - {w}")
        if len(warnings) > 10:
            print(
                f"  ... and {len(warnings) - 10} more (see 'warnings' in output file)"
            )
    print("-" * 60)
    print(f"Written: {output_file.resolve()}")
    print("=" * 60 + "\n")


if __name__ == "__main__":
    main()
